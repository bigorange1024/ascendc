#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""合稿衔接片：10 页。白底；冷色主调、暖色强调。

硬约束：示意图一律用 python-pptx 原生形状 + TextFrame，禁止 add_picture。
内容来源：
  - 第2页：用户样板（挑战令主题 + 四约束）
  - 第3–4页：thirdparty/ntt_onnx/docs/summary/全流程工程记述（§7.7–7.8 人机协作）
  - 第5–8页：docs/research/从已验证能力到合法派生-…教材草案（问题/方法/效果）
再生：python3 docs/reports/gen_ai_bounded_search_briefing_pptx.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from lxml import etree

OUT = "/workspace/docs/reports/汇报_有界推进与绿灯底座_合稿衔接片.pptx"
W, H = Inches(13.333), Inches(7.5)

C_BG = RGBColor(0xFF, 0xFF, 0xFF)
C_INK = RGBColor(0x1A, 0x2A, 0x3A)
C_MUTED = RGBColor(0x3D, 0x4F, 0x5F)
C_ACCENT = RGBColor(0x1E, 0x5A, 0x7A)
C_DEEP = RGBColor(0x14, 0x32, 0x52)
C_OK = RGBColor(0x1F, 0x6B, 0x5E)
C_OK_SOFT = RGBColor(0xE6, 0xF3, 0xEF)
C_WARM = RGBColor(0xC4, 0x5C, 0x26)
C_WARM_SOFT = RGBColor(0xFD, 0xF0, 0xE6)
C_SOFT = RGBColor(0xE8, 0xF0, 0xF6)
C_SOFT2 = RGBColor(0xD5, 0xE4, 0xEF)
C_GRAY = RGBColor(0x90, 0xA0, 0xAC)
C_GRAY_SOFT = RGBColor(0xF2, 0xF4, 0xF6)
C_BAND = RGBColor(0xEE, 0xF3, 0xF7)
C_HEAD = RGBColor(0xE3, 0xED, 0xF5)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Microsoft YaHei"


def _run(run, text, size=18, bold=False, color=C_INK):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    if ea is None:
        ea = etree.SubElement(rPr, "{http://schemas.openxmlformats.org/drawingml/2006/main}ea")
    ea.set("typeface", FONT)


def fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def stroke(shape, color, width=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def set_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill(bg, C_BG)
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def tb(slide, left, top, width, height, text, size=18, bold=False, color=C_INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    _run(r, text, size=size, bold=bold, color=color)
    return box


def rich(slide, left, top, width, height, paragraphs, default_size=17):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in paragraphs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(3)
        if isinstance(item, list):
            for run_item in item:
                t = run_item[0]
                sz = run_item[1] if len(run_item) > 1 else default_size
                bd = run_item[2] if len(run_item) > 2 else False
                cl = run_item[3] if len(run_item) > 3 else C_INK
                r = p.add_run()
                _run(r, t, size=sz, bold=bd, color=cl)
        elif isinstance(item, str):
            r = p.add_run()
            _run(r, item, size=default_size, bold=False, color=C_INK)
        else:
            t = item[0]
            sz = item[1] if len(item) > 1 else default_size
            bd = item[2] if len(item) > 2 else False
            cl = item[3] if len(item) > 3 else C_INK
            r = p.add_run()
            _run(r, t, size=sz, bold=bd, color=cl)
    return box


def shape_text(shape, lines, size=14, bold=False, color=C_INK, align=PP_ALIGN.CENTER, margin_top=6):
    """lines: str | [str|(t,sz,bd,cl) | [(t,sz,bd,cl),...混色行] ]"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    if isinstance(lines, str):
        parts_per_line = [[(ln, size, bold, color)] for ln in lines.split("\n")]
    else:
        parts_per_line = []
        for item in lines:
            if isinstance(item, str):
                parts_per_line.append([(item, size, bold, color)])
            elif isinstance(item, list):
                # 同行多 run
                runs = []
                for sub in item:
                    if isinstance(sub, str):
                        runs.append((sub, size, bold, color))
                    else:
                        t = sub[0]
                        sz = sub[1] if len(sub) > 1 else size
                        bd = sub[2] if len(sub) > 2 else bold
                        cl = sub[3] if len(sub) > 3 else color
                        runs.append((t, sz, bd, cl))
                parts_per_line.append(runs)
            else:
                t = item[0]
                sz = item[1] if len(item) > 1 else size
                bd = item[2] if len(item) > 2 else bold
                cl = item[3] if len(item) > 3 else color
                parts_per_line.append([(t, sz, bd, cl)])
    for i, runs in enumerate(parts_per_line):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(2)
        for t, sz, bd, cl in runs:
            run = p.add_run()
            _run(run, t, size=sz, bold=bd, color=cl)
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(margin_top)
    tf.margin_bottom = Pt(4)


def rounded(slide, left, top, width, height, lines, fill_c=C_SOFT, line_c=C_ACCENT,
            text_c=C_INK, size=14, bold=False, align=PP_ALIGN.CENTER, margin_top=8, line_w=1.25):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.adjustments[0] = 0.06
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_c
    stroke(sh, line_c, line_w)
    shape_text(sh, lines, size=size, bold=bold, color=text_c, align=align, margin_top=margin_top)
    return sh


def arrow_right(slide, left, top, width, height, color=C_ACCENT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    fill(sh, color)
    return sh


def arrow_down(slide, left, top, width, height, color=C_WARM):
    sh = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    fill(sh, color)
    return sh


def title_theme(slide, text):
    tb(slide, Inches(0.4), Inches(0.14), Inches(12.5), Inches(0.38), text, size=20, bold=True, color=C_DEEP)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.52), Inches(0.85), Inches(0.04))
    fill(bar, C_ACCENT)


def page_summary(slide, text, page, total):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(6.9), Inches(12.6), Inches(0.4))
    fill(band, C_BAND)
    tb(slide, Inches(0.5), Inches(6.96), Inches(11.2), Inches(0.3), f"总结：{text}", size=13, bold=True, color=C_DEEP)
    tb(slide, Inches(11.85), Inches(6.98), Inches(0.95), Inches(0.28), f"{page}/{total}", size=12, color=C_MUTED, align=PP_ALIGN.RIGHT)


def section_label(slide, left, top, width, height, text, fill_c=C_HEAD, line_c=C_ACCENT, text_c=C_DEEP):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.adjustments[0] = 0.08
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_c
    stroke(sh, line_c, 1.1)
    shape_text(sh, [(text, 14, True, text_c)], size=14, bold=True, color=text_c, margin_top=8)
    return sh


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # ========== 1 封面（暂不改结构，仅轻量）==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    panel = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), H)
    fill(panel, C_DEEP)
    tb(s, Inches(0.55), Inches(1.35), Inches(12), Inches(0.35),
       "AI 辅助硬科技研发 · 合稿材料", size=16, bold=True, color=C_ACCENT)
    tb(s, Inches(0.55), Inches(1.85), Inches(12), Inches(0.7),
       "有界推进与绿灯底座", size=40, bold=True, color=C_DEEP)
    rich(s, Inches(0.55), Inches(2.7), Inches(12), Inches(1.55), [
        ("这场汇报讲一条能力线，分两段。", 17),
        ("上篇：过去那个项目，在四重约束下怎么管住 AI。", 17),
        ("下篇：同一条线怎么升级成「绿灯底座 + 缺项里补」，冲固定可核对的终点。", 17),
        ("主语是 Agent 工程能力；ML-KEM 是检验床，不是算法仓库路演。", 17),
    ])
    rounded(s, Inches(0.7), Inches(4.55), Inches(4.9), Inches(1.15),
            [("上篇", 15, True, C_DEEP), ("过去：四约束下管住 AI", 15, False, C_INK)],
            fill_c=C_SOFT, line_c=C_DEEP, margin_top=16)
    arrow_right(s, Inches(5.85), Inches(4.9), Inches(1.15), Inches(0.45), C_WARM)
    tb(s, Inches(5.55), Inches(4.45), Inches(1.8), Inches(0.35),
       "同一条能力线", size=12, bold=True, color=C_WARM, align=PP_ALIGN.CENTER)
    rounded(s, Inches(7.25), Inches(4.55), Inches(5.3), Inches(1.15),
            [("下篇", 15, True, C_ACCENT), ("今天：底座 × 缺项冲固定终点", 15, False, C_INK)],
            fill_c=C_SOFT, line_c=C_ACCENT, margin_top=16)
    tb(s, Inches(11.85), Inches(6.95), Inches(1.0), Inches(0.3), "1/10", size=12, color=C_MUTED, align=PP_ALIGN.RIGHT)

    # ========== 2 用户样板：挑战令主题 + 四约束 ==========
    # 口径对齐用户截图写法（挑战令主题之一…）
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_theme(s, "挑战令主题之一：基于存量旧型号 AI Core 构造高性能 PQC 关键算子实现")

    section_label(s, Inches(0.4), Inches(0.62), Inches(1.35), Inches(0.36), "要做什么")
    rich(s, Inches(1.9), Inches(0.62), Inches(11.0), Inches(1.05), [
        [
            ("在目标芯片的 ", 14, False, C_INK),
            ("AI Core", 14, True, C_ACCENT),
            (" 上，把后量子密码密钥生成里约定好的 ", 14, False, C_INK),
            ("NTT 计算", 14, True, C_ACCENT),
            (" 做出来。", 14, False, C_INK),
        ], [
            ("硬验收：", 14, True, C_DEEP),
            ("实机真能跑", 14, True, C_WARM),
            ("；输出与参考程序 ", 14, False, C_INK),
            ("逐项对上", 14, True, C_WARM),
            ("。技术路线不规定。", 14, False, C_INK),
        ], [
            ("推进节奏：C 语言基准代码 → 方案设计 → 探路（定主线）→ AI Core 实现（基准→批量→融合算子）→ 测试+交付。", 13, False, C_MUTED),
        ],
    ], default_size=14)

    section_label(s, Inches(0.4), Inches(1.75), Inches(1.35), Inches(0.36), "为什么难")
    rounded(s, Inches(1.9), Inches(1.72), Inches(10.95), Inches(0.48),
            [[
                ("核心挑战：", 14, True, C_DEEP),
                ("路线不确定 × 平台不成熟 × 语言不常见 × 设备不直调", 14, True, C_WARM),
            ]],
            fill_c=C_SOFT, line_c=C_ACCENT, align=PP_ALIGN.LEFT, margin_top=8)

    cards = [
        ("① 路线不确定",
         "造成的困难：研究类项目失败尝试特别多。新旧代码叠放，主线未钉死时探索成本极高、状态难对齐。",
         "AI 跑偏：很难识别有价值的路线，容易把死路当捷径，在错误方向原地打转。"),
        ("② 平台不成熟",
         "造成的困难：硬件小众，外网几乎无现成解法，只能靠本机实验把路走通。",
         "AI 跑偏：没有全面的数据库，用通识 / 公网经验硬套，推断常错却很自信。"),
        ("③ 语言不常见",
         "造成的困难：昇腾专用编程，API 是黑盒（还有 bug）。调试链路长；仿真结果经常与真机不同。",
         "AI 跑偏：模型缺乏可靠的编译先验知识。"),
        ("④ 设备不直调",
         "造成的困难：AI 够不着实机，终测须人上机；日志须人工脱敏，AI 无法远程代终审。",
         "AI 跑偏：把仿真通过当成项目结项；基于错误前提推理。"),
    ]
    positions = [(0.4, 2.4), (6.85, 2.4), (0.4, 4.45), (6.85, 4.45)]
    for (title, hard, fail), (x, y) in zip(cards, positions):
        rounded(s, Inches(x), Inches(y), Inches(6.0), Inches(1.8),
                [(title, 15, True, C_DEEP),
                 (hard, 13, False, C_INK),
                 (fail, 13, True, C_WARM)],
                fill_c=C_SOFT, line_c=C_ACCENT, align=PP_ALIGN.LEFT, margin_top=10)
    page_summary(s, "任务有固定验收，但其余高度不确定；仅靠 AI 无法解决实际开发问题。", 2, 10)

    # ========== 3 如何接住（全流程工程记述 §7.7）==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_theme(s, "主题：当时怎么接住——人掌舵，Agent 在边界内做工")
    tb(s, Inches(0.4), Inches(0.58), Inches(12.5), Inches(0.28),
       "摘自《全流程工程记述》§7.7 人机协作：偏研究型工程中的 Agent 经验", size=12, color=C_MUTED)

    rich(s, Inches(0.4), Inches(0.9), Inches(12.5), Inches(0.7), [
        [
            ("在四约束叠加条件下用 Agent，若缺仓库状态锚点与可执行验收，极易出现：", 14, False, C_INK),
        ], [
            ("「对话里已结项、仓库里仍 pending」", 14, True, C_WARM),
            ("；", 14, False, C_INK),
            ("「声称实机通过却无 cmp 记录」", 14, True, C_WARM),
            ("。", 14, False, C_INK),
        ], [
            ("定调：", 14, True, C_DEEP),
            ("人掌舵、Agent 在边界内做工制品", 14, True, C_WARM),
            ("——不是「Agent 承包项目」。", 14, False, C_INK),
        ],
    ], default_size=14)

    # 闭环箭头
    rounded(s, Inches(0.4), Inches(1.85), Inches(2.7), Inches(0.95),
            [("人写入", 13, True, C_DEEP), ("目标·范围·验收", 12, False, C_INK)],
            fill_c=C_WARM_SOFT, line_c=C_WARM, margin_top=14, line_w=1.5)
    arrow_right(s, Inches(3.2), Inches(2.1), Inches(0.45), Inches(0.35), C_ACCENT)
    rounded(s, Inches(3.75), Inches(1.85), Inches(2.9), Inches(0.95),
            [("Agent 做工", 13, True, C_ACCENT), ("契约内改制品+回传输出", 12, False, C_INK)],
            fill_c=C_SOFT, line_c=C_ACCENT, margin_top=14)
    arrow_right(s, Inches(6.75), Inches(2.1), Inches(0.45), Inches(0.35), C_ACCENT)
    rounded(s, Inches(7.3), Inches(1.85), Inches(2.7), Inches(0.95),
            [("人终审", 13, True, C_WARM), ("独占环境验证", 12, False, C_INK)],
            fill_c=C_WARM_SOFT, line_c=C_WARM, margin_top=14, line_w=1.5)
    arrow_right(s, Inches(10.1), Inches(2.1), Inches(0.45), Inches(0.35), C_OK)
    rounded(s, Inches(10.65), Inches(1.85), Inches(2.2), Inches(0.95),
            [("写回账本", 13, True, C_OK), ("单一事实源", 12, False, C_INK)],
            fill_c=C_OK_SOFT, line_c=C_OK, margin_top=14)

    # 两栏：人 / Agent
    rounded(s, Inches(0.4), Inches(3.05), Inches(6.15), Inches(2.55),
            [("人 · 主要负责（表38）", 14, True, C_WARM),
             ("意义与边界：问题、验收口径、已关闭路线", 13, False, C_INK),
             ("取舍：路线采纳 / 终止决策", 13, False, C_INK),
             ("独占环境：专有编译器、目标机、物理设备", 13, False, C_INK),
             ("认定完成：凭可复验证据裁定；范围变更签字", 13, False, C_INK)],
            fill_c=C_WARM_SOFT, line_c=C_WARM, align=PP_ALIGN.LEFT, margin_top=10, line_w=1.5)
    rounded(s, Inches(6.75), Inches(3.05), Inches(6.15), Inches(2.55),
            [("Agent · 主要负责（表38）", 14, True, C_ACCENT),
             ("制品操作：检索、改代码/脚本/文档草稿", 13, False, C_INK),
             ("可代劳执行：授权范围内构建/测试并回传完整输出", 13, False, C_INK),
             ("不宜：对外承诺范围/工期；无 log 断言实机已过", 13, True, C_WARM),
             ("不宜：擅自扩大范围；复活已关闭路线", 13, True, C_WARM)],
            fill_c=C_SOFT, line_c=C_ACCENT, align=PP_ALIGN.LEFT, margin_top=10)
    page_summary(s, "交付凭证据不靠叙述；换路权与终审权在人；Agent 干可核对中间层。", 3, 10)

    # ========== 4 可带走习惯（全流程 §7.8）+ 关键转折 ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_theme(s, "主题：可带走什么——协作习惯，不是某一版源码")
    tb(s, Inches(0.4), Inches(0.58), Inches(12.5), Inches(0.28),
       "摘自《全流程工程记述》§7.8 可带到其他项目的协作习惯；关键转折见第二章", size=12, color=C_MUTED)

    # 关键转折条
    rounded(s, Inches(0.4), Inches(0.95), Inches(12.5), Inches(0.7),
            [[
                ("关键转折：", 14, True, C_WARM),
                ("CCE 直连闭环慢、可验证性弱 → 人拍板改走「计算图 → 模型转换 → 实机」为主线", 13, False, C_INK),
            ], [
                ("验证链定型：", 13, True, C_DEEP),
                ("主机 C 参考 → 上板前核对 → 实机验收（过关 = 输出全等）", 13, False, C_INK),
            ]],
            fill_c=C_WARM_SOFT, line_c=C_WARM, align=PP_ALIGN.LEFT, margin_top=6, line_w=1.5)

    habits = [
        ("仓库写清范围", "账本：结项/主线/已关闭表；禁止只用聊天改范围"),
        ("做完能核对", "完成 = 可执行检查（cmp / 退出码 / 完整 log）"),
        ("验证分层", "主机 ≠ 硅前 ≠ 实机；子图过 ≠ 整图过"),
        ("试制进交付须人批", "AI 不宣布结项；人跑通 pipeline 后签字"),
        ("中断写清接续", "已完成(证据)·当前判断·下一步一条·勿重做"),
        ("任务单写死边界", "目标·必读·允许改·禁止·DoD·回报未做项"),
    ]
    for i, (a, b) in enumerate(habits):
        col, row = i % 3, i // 3
        left = 0.4 + col * 4.25
        top = 1.9 + row * 1.85
        rounded(s, Inches(left), Inches(top), Inches(4.05), Inches(1.6),
                [(a, 15, True, C_DEEP), (b, 13, False, C_INK)],
                fill_c=C_SOFT, line_c=C_ACCENT, align=PP_ALIGN.LEFT, margin_top=18)
    page_summary(s, "习惯写进仓库才能换会话接续；这四条种子直接长成下篇「底座 / 缺项 / 交叉」。", 4, 10)

    # ========== 5 今天：问题与主张（教材 §3.1）——密文，勿大方框 ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_theme(s, "主题：今天难在哪、主张什么")
    tb(s, Inches(0.4), Inches(0.56), Inches(12.5), Inches(0.25),
       "摘自教材草案 §3.1", size=11, color=C_MUTED)

    section_label(s, Inches(0.4), Inches(0.88), Inches(1.2), Inches(0.34), "引子")
    rich(s, Inches(1.75), Inches(0.85), Inches(11.1), Inches(1.15), [
        [
            ("correctness 有了，任务就算完成了吗？", 14, True, C_WARM),
            (" 仓库底层能力大致齐了，提示很轻：让 Agent 自行读标准、把封装/解封装「做出来」。", 13, False, C_INK),
        ], [
            ("结果能对上，像是做完了；摊开却是：", 13, False, C_INK),
            ("只保结果对、几乎不为架构设计", 13, True, C_WARM),
            ("；搜索不设界、反复碰壁；", 13, False, C_INK),
            ("短时间耗尽约一个月配额", 13, True, C_WARM),
            ("；还沉淀不正确的「经验」。", 13, False, C_INK),
        ], [
            ("伪完成：", 13, True, C_DEEP),
            ("结果碰巧对，过程发散、形态不可用、钱与时间被烧掉。", 13, True, C_WARM),
        ],
    ], default_size=13)

    section_label(s, Inches(0.4), Inches(2.15), Inches(1.2), Inches(0.34), "问题实质")
    rich(s, Inches(1.75), Inches(2.12), Inches(11.1), Inches(1.2), [
        [
            ("开放生成被当成了合法派生。", 14, True, C_WARM),
            (" 菜单里有、局部测通、文本上看起来像解，就被当成「整体已经合法、可以交付」。", 13, False, C_INK),
        ], [
            ("工程语境里的「幻觉」，多半是", 13, False, C_INK),
            ("无效派生", 13, True, C_WARM),
            ("：把未获证节点当可引用能力，把接缝跳过去。", 13, False, C_INK),
        ], [
            ("要盯住的不是「会不会写出对拍代码」，而是：", 13, False, C_INK),
            ("如何阻止把无界生成当作完成预研的合法途径", 13, True, C_DEEP),
            ("。", 13, False, C_INK),
        ],
    ], default_size=13)

    section_label(s, Inches(0.4), Inches(3.5), Inches(1.2), Inches(0.34), "方法实质")
    rich(s, Inches(1.75), Inches(3.45), Inches(11.1), Inches(1.35), [
        [
            ("把无界生成收成：", 13, False, C_INK),
            ("已认证闭包上的受控搜索", 14, True, C_ACCENT),
            ("。别当自由写手，当证书闭包上的搜索器。", 13, False, C_INK),
        ], [
            ("① 已验证能力：", 13, True, C_DEEP),
            ("现在还能引用谁；  ", 13, False, C_INK),
            ("② 禁止项：", 13, True, C_WARM),
            ("明确不能再当模板的路；  ", 13, False, C_INK),
            ("③ 合法拼装：", 13, True, C_ACCENT),
            ("缺的先补齐并验过，再往上引用。", 13, False, C_INK),
        ], [
            ("「有限」指：相对当前能力集/禁止项/拼装规则，下一步可枚举、可检查；集合可随新能力获证长大。", 13, False, C_MUTED),
        ],
    ], default_size=13)

    section_label(s, Inches(0.4), Inches(5.0), Inches(1.2), Inches(0.34), "本场主张", fill_c=C_WARM_SOFT, line_c=C_WARM, text_c=C_WARM)
    rich(s, Inches(1.75), Inches(4.95), Inches(11.1), Inches(1.2), [
        [
            ("只冲答案事先钉死、可核对的固定终点。", 14, True, C_WARM),
            (" 开放发现可与别人互补，但不是本场主语。", 13, False, C_INK),
        ], [
            ("额度有限、多模型水平不一 → ", 13, False, C_INK),
            ("换谁上场都要能干", 13, True, C_DEEP),
            ("；搜索边界必须可检查。", 13, False, C_INK),
        ],
    ], default_size=13)
    page_summary(s, "主张：把开放生成收成沿已验证能力的可检查受控搜索。", 5, 10)

    # ========== 6 底座×缺项——密文分段 ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_theme(s, "主题：绿灯底座从哪来，缺项怎么补")
    tb(s, Inches(0.4), Inches(0.56), Inches(12.5), Inches(0.25),
       "摘自教材草案：架构直觉 + 写码许可（先交闭包表，再写码）", size=11, color=C_MUTED)

    section_label(s, Inches(0.4), Inches(0.88), Inches(1.35), Inches(0.34), "底座从哪来")
    rich(s, Inches(1.9), Inches(0.85), Inches(11.0), Inches(1.15), [
        [
            ("人先把能用的沉下来：", 13, False, C_INK),
            ("编程笔记、跑通的原型、带着 Agent 走过的长轨迹", 13, True, C_ACCENT),
            ("（如 KeyGen）。", 13, False, C_INK),
        ], [
            ("写进索引、验过的才算绿灯底座 Γ；", 13, False, C_INK),
            ("聊天里「好像会了」不算", 13, True, C_WARM),
            ("。活跃菜单是线索，", 13, False, C_INK),
            ("不等于合法拼装说明书", 13, True, C_WARM),
            ("。", 13, False, C_INK),
        ],
    ], default_size=13)

    section_label(s, Inches(0.4), Inches(2.15), Inches(1.35), Inches(0.34), "缺项怎么定")
    rich(s, Inches(1.9), Inches(2.12), Inches(11.0), Inches(1.15), [
        [
            ("做成目标 T，理论上要齐哪些能力，减去已经绿灯的，剩下就是缺项。", 13, False, C_INK),
        ], [
            ("只许在缺项里补；", 13, True, C_WARM),
            ("补出来的还要再验。两块都绿，", 13, False, C_INK),
            ("接缝照样单独过", 13, True, C_WARM),
            ("（组合不免费）。", 13, False, C_INK),
        ], [
            ("死路 / frozen / 已关闭路线 → ", 13, False, C_INK),
            ("禁止项，不许当模板", 13, True, C_WARM),
            ("。", 13, False, C_INK),
        ],
    ], default_size=13)

    section_label(s, Inches(0.4), Inches(3.45), Inches(1.35), Inches(0.34), "三件对象")
    rich(s, Inches(1.9), Inches(3.42), Inches(11.0), Inches(1.35), [
        [
            ("① 已验证能力：", 13, True, C_DEEP),
            ("现在还能引用谁。  ", 13, False, C_INK),
            ("② 禁止项：", 13, True, C_WARM),
            ("明确不能再走的路。  ", 13, False, C_INK),
            ("③ 合法拼装 / 出口交叉：", 13, True, C_ACCENT),
            ("怎样衔接才算一步合法推进。", 13, False, C_INK),
        ], [
            ("门禁一句话：", 13, True, C_DEEP),
            ("先交闭包表（目标、已获证、缺项、接缝、禁止相交），再写码。", 13, True, C_WARM),
        ], [
            ("customspec / STATUS 描述「怎么实现」；闭包表回答「凭什么现在允许开始实现」。", 12, False, C_MUTED),
        ],
    ], default_size=13)

    section_label(s, Inches(0.4), Inches(5.0), Inches(1.35), Inches(0.34), "与上篇关系", fill_c=C_WARM_SOFT, line_c=C_WARM, text_c=C_WARM)
    rich(s, Inches(1.9), Inches(4.95), Inches(11.0), Inches(1.15), [
        [
            ("上篇习惯直接长成这里：", 13, False, C_INK),
            ("账本→底座+禁止项", 13, True, C_DEEP),
            ("；", 13, False, C_INK),
            ("做完定义→出口可核对/交叉", 13, True, C_DEEP),
            ("；", 13, False, C_INK),
            ("换会话能接→索引可交接", 13, True, C_DEEP),
            ("；", 13, False, C_INK),
            ("任务单→写码边界", 13, True, C_DEEP),
            ("。", 13, False, C_INK),
        ], [
            ("不是另起炉灶，是同一条能力线往上长。", 13, True, C_WARM),
        ],
    ], default_size=13)
    page_summary(s, "搜索收到缺项里；禁止项挡死路；先表后码。", 6, 10)

    # ========== 7 效果：A/B/C 对照图（correctness × 手调）——立刻可用 ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_theme(s, "主题：效果——对照 correctness，对照手调精修")
    tb(s, Inches(0.4), Inches(0.54), Inches(12.5), Inches(0.24),
       "检验床：Ascend ML-KEM（Decaps 试金石 + 768 体感）。棒图为体感示意，不作严格倍数账单。", size=11, color=C_MUTED)

    # —— 三柱对照图：A / B / C ——
    # A
    rounded(s, Inches(0.35), Inches(0.85), Inches(4.0), Inches(2.55),
            [("A  无方法论", 16, True, C_DEEP),
             ("correctness 开放探索", 13, True, C_ACCENT),
             ("只求结果对 · 搜索不设界", 12, False, C_INK),
             ("正确性：能对上", 12, False, C_INK),
             ("性能/形态：常不可用", 12, False, C_INK),
             ("过程成本：很高", 13, True, C_WARM),
             ("判决：更费（反面教材）", 13, True, C_WARM)],
            fill_c=C_SOFT, line_c=C_ACCENT, align=PP_ALIGN.LEFT, margin_top=8)
    # B 主叙事（暖边）
    rounded(s, Inches(4.55), Inches(0.85), Inches(4.2), Inches(2.55),
            [("B  方法论  ★主叙事", 16, True, C_WARM),
             ("底座 + 缺项 + 出口交叉", 13, True, C_OK),
             ("先交表再写码 · 禁抄死路", 12, False, C_INK),
             ("正确性：与 C 同阶", 12, False, C_INK),
             ("性能：与 C 同量级（未宣称更快）", 12, False, C_INK),
             ("过程成本：明显收小", 13, True, C_OK),
             ("判决：相对 A 强成功", 13, True, C_OK)],
            fill_c=C_OK_SOFT, line_c=C_WARM, align=PP_ALIGN.LEFT, margin_top=8, line_w=2.0)
    # C
    rounded(s, Inches(8.95), Inches(0.85), Inches(3.95), Inches(2.55),
            [("C  手调精修", 16, True, C_DEEP),
             ("人工长期优化交付树", 13, True, C_DEEP),
             ("基线 · 工程主线默认", 12, False, C_INK),
             ("正确性：交付级", 12, False, C_INK),
             ("性能：人工精修定标", 12, False, C_INK),
             ("过程：人带着走", 12, False, C_INK),
             ("对照结论：B ≈ C", 13, True, C_DEEP)],
            fill_c=C_SOFT, line_c=C_DEEP, align=PP_ALIGN.LEFT, margin_top=8)

    # —— 关系箭头说明（短）——
    tb(s, Inches(0.4), Inches(3.5), Inches(12.5), Inches(0.28),
       "读图：相对 A → B 赢在过程与结构；相对 C → B 正确性/tick 同量级，不宣称性能胜出。方法论代替不了人工优化主线。",
       size=12, bold=False, color=C_MUTED)

    # —— 过程负担横向棒图 ——
    tb(s, Inches(0.4), Inches(3.85), Inches(6.2), Inches(0.28),
       "过程负担体感（示意）", size=13, bold=True, color=C_DEEP)

    # A bar
    tb(s, Inches(0.4), Inches(4.2), Inches(2.0), Inches(0.28), "A 开放 correctness", size=11, color=C_MUTED)
    ba = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(4.2), Inches(5.8), Inches(0.32))
    ba.adjustments[0] = 0.1
    fill(ba, C_ACCENT)
    tb(s, Inches(8.4), Inches(4.2), Inches(1.5), Inches(0.28), "很高", size=12, bold=True, color=C_ACCENT)

    # B bar
    tb(s, Inches(0.4), Inches(4.65), Inches(2.0), Inches(0.28), "B 有门禁方法论", size=11, color=C_MUTED)
    bb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(4.65), Inches(1.6), Inches(0.32))
    bb.adjustments[0] = 0.1
    fill(bb, C_OK)
    tb(s, Inches(4.2), Inches(4.65), Inches(2.0), Inches(0.28), "收小", size=12, bold=True, color=C_OK)

    # C bar (reference, medium)
    tb(s, Inches(0.4), Inches(5.1), Inches(2.0), Inches(0.28), "C 手调精修（参照）", size=11, color=C_MUTED)
    bc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.1), Inches(2.0), Inches(0.32))
    bc.adjustments[0] = 0.1
    fill(bc, C_DEEP)
    tb(s, Inches(4.65), Inches(5.1), Inches(3.5), Inches(0.28), "人带着走（基线）", size=12, bold=True, color=C_DEEP)

    # —— 右侧：768 体感小卡 ——
    rounded(s, Inches(8.3), Inches(3.85), Inches(4.55), Inches(2.0),
            [("768 体感（不作倍数宣称）", 13, True, C_WARM),
             ("以前 correctness 开放搜：", 11, False, C_INK),
             ("约 1 天墙钟 · 月配额大头", 12, True, C_WARM),
             ("有门禁一夜：", 11, False, C_INK),
             ("配额 +~14pt · 墙钟 <4h", 12, True, C_OK),
             ("省的是无效派生，非魔法提速", 11, False, C_MUTED)],
            fill_c=C_WARM_SOFT, line_c=C_WARM, align=PP_ALIGN.LEFT, margin_top=6, line_w=1.5)

    page_summary(s, "B≈手调精修；相对 correctness 更省。主叙事是 B。", 7, 10)

    # ========== 8 合稿——密文 ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_theme(s, "主题：跟别人的材料怎么接得上")

    section_label(s, Inches(0.4), Inches(0.7), Inches(1.35), Inches(0.34), "合稿怕什么")
    rich(s, Inches(1.9), Inches(0.68), Inches(11.0), Inches(1.0), [
        [
            ("最怕各说各话。", 14, True, C_WARM),
            (" 别人若在某场景把某一招打磨得比现成方案更强——那正好是我们要的「更好积木」。", 13, False, C_INK),
        ], [
            ("我们", 13, False, C_INK),
            ("不做开放式搜灵感", 13, True, C_MUTED),
            ("。把拿到手的积木登记进绿灯底座，再按缺项清单组装到事先钉死的目标。", 13, False, C_INK),
        ],
    ], default_size=13)

    section_label(s, Inches(0.4), Inches(1.9), Inches(1.35), Inches(0.34), "交接一句", fill_c=C_WARM_SOFT, line_c=C_WARM, text_c=C_WARM)
    rich(s, Inches(1.9), Inches(1.88), Inches(11.0), Inches(0.85), [
        [
            ("他们输出更优能力点，我们装进底座并走到终点。", 15, True, C_WARM),
        ], [
            ("底座变好，终点实现跟着变好——与「DAG / 积木迭代创新」互补，", 13, False, C_INK),
            ("不抢开放发现主语", 13, True, C_DEEP),
            ("。", 13, False, C_INK),
        ],
    ], default_size=13)

    section_label(s, Inches(0.4), Inches(2.95), Inches(1.35), Inches(0.34), "分工怎么挂")
    rich(s, Inches(1.9), Inches(2.92), Inches(11.0), Inches(1.5), [
        [
            ("对方：", 13, True, C_ACCENT),
            ("在场景里反复打磨，输出比现成方案更强的能力点 / 积木。", 13, False, C_INK),
        ], [
            ("我方：", 13, True, C_DEEP),
            ("注入 Γ（登记、获证）→ 按缺项组装到固定可核对终点 T。", 13, False, C_INK),
        ], [
            ("前提：", 13, True, C_WARM),
            ("上篇那套习惯（账本、做完定义、可交接、任务边界）——否则接不住别人的积木，只会重新发散。", 13, False, C_INK),
        ],
    ], default_size=13)

    section_label(s, Inches(0.4), Inches(4.6), Inches(1.35), Inches(0.34), "领导带走")
    rich(s, Inches(1.9), Inches(4.55), Inches(11.0), Inches(1.35), [
        [
            ("合稿不是两套故事硬贴，而是一条接口：", 13, False, C_INK),
            ("更优积木 → 绿灯底座 → 固定终点", 14, True, C_WARM),
            ("。", 13, False, C_INK),
        ], [
            ("我们负责把不确定性收到可检查边界里，并把别人的增量装进能交付的底座。", 13, False, C_INK),
        ],
    ], default_size=13)
    page_summary(s, "挂钩：更优积木 → 绿灯底座 → 固定终点。", 8, 10)

    # ========== 9 四阶段——密文，非空心流程图 ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_theme(s, "主题：过去项目怎么推——四阶段与验证链")
    tb(s, Inches(0.4), Inches(0.56), Inches(12.5), Inches(0.25),
       "摘自《全流程工程记述》§1.2 / §1.3.1 / 第二章关键转折", size=11, color=C_MUTED)

    section_label(s, Inches(0.4), Inches(0.88), Inches(1.35), Inches(0.34), "四阶段")
    rich(s, Inches(1.9), Inches(0.85), Inches(11.0), Inches(1.55), [
        [
            ("① 准备：", 13, True, C_DEEP),
            ("澄清交付边界与验收口径；主机侧建立与语义一致的参考实现及可对拍数据。", 13, False, C_INK),
        ], [
            ("② 探路：", 13, True, C_DEEP),
            ("试跑「参考→实机」多条路径，比可验证性与迭代效率，选定交付主线。", 13, False, C_INK),
        ], [
            ("③ 实施：", 13, True, C_DEEP),
            ("在主路径上完成开发、分层验证与实机对拍（基准→批量→合成）。", 13, False, C_INK),
        ], [
            ("④ 收口：", 13, True, C_DEEP),
            ("与客户确认部署范围；整理可复现材料与索引快照。", 13, False, C_INK),
        ],
    ], default_size=13)

    section_label(s, Inches(0.4), Inches(2.6), Inches(1.35), Inches(0.34), "关键转折", fill_c=C_WARM_SOFT, line_c=C_WARM, text_c=C_WARM)
    rich(s, Inches(1.9), Inches(2.55), Inches(11.0), Inches(1.25), [
        [
            ("CCE 直连试跑：", 13, False, C_INK),
            ("闭环慢、可验证性弱、难维持每日核对", 13, True, C_WARM),
            (" → 保留为对照，", 13, False, C_INK),
            ("不再作正式交付主线", 13, True, C_WARM),
            ("。", 13, False, C_INK),
        ], [
            ("人拍板改走：", 13, True, C_DEEP),
            ("计算图 → 模型转换 → 实机 AI Core", 13, True, C_ACCENT),
            ("。硅前可一日多次迭代，问题可分层定位。", 13, False, C_INK),
        ], [
            ("换路权在人，不在 Agent。", 13, True, C_WARM),
        ],
    ], default_size=13)

    section_label(s, Inches(0.4), Inches(4.0), Inches(1.35), Inches(0.34), "验证链")
    rich(s, Inches(1.9), Inches(3.95), Inches(11.0), Inches(1.5), [
        [
            ("① 参考基准：", 13, True, C_DEEP),
            ("主机 C 与标准输入输出。  ", 13, False, C_INK),
            ("② 上板前核对：", 13, True, C_DEEP),
            ("同输入跑待交付件，与参考比。  ", 13, False, C_INK),
            ("③ 实机验收：", 13, True, C_DEEP),
            ("AI Core 输出须与参考/上板前一致。", 13, False, C_INK),
        ], [
            ("过关口径：", 13, True, C_WARM),
            ("输出全等 / max_abs_diff=0", 13, True, C_WARM),
            ("——不是「代码读起来合理」，不是「对话里说通过了」。", 13, False, C_INK),
        ], [
            ("须写清测到哪一层：硅前通过 ≠ 实机已验收；子图通过 ≠ 整图通过。", 13, False, C_MUTED),
        ],
    ], default_size=13)
    page_summary(s, "四阶段推进；CCE 推不动由人换路；验收分层、凭证据。", 9, 10)

    # ========== 10 收束——密文，无空心能力线大框 ==========
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    title_theme(s, "主题：今天带回去哪几句")

    section_label(s, Inches(0.4), Inches(0.75), Inches(1.0), Inches(0.34), "一")
    rich(s, Inches(1.55), Inches(0.72), Inches(11.3), Inches(0.85), [
        [
            ("四约束下可以把 AI 管住", 15, True, C_WARM),
            ("——人掌舵、证据落盘、做完能核对、换会话能接。", 14, False, C_INK),
        ], [
            ("资产是规矩与习惯，不是某一版源码。", 13, False, C_MUTED),
        ],
    ], default_size=14)

    section_label(s, Inches(0.4), Inches(1.75), Inches(1.0), Inches(0.34), "二")
    rich(s, Inches(1.55), Inches(1.72), Inches(11.3), Inches(0.95), [
        [
            ("同一条能力线升级为：", 14, False, C_INK),
            ("绿灯底座 × 缺项里补", 15, True, C_WARM),
            ("，只冲固定可核对终点。", 14, False, C_INK),
        ], [
            ("把开放生成收成可检查的受控搜索；先表后码。", 13, False, C_MUTED),
        ],
    ], default_size=14)

    section_label(s, Inches(0.4), Inches(2.85), Inches(1.0), Inches(0.34), "三")
    rich(s, Inches(1.55), Inches(2.82), Inches(11.3), Inches(0.95), [
        [
            ("合稿点：", 14, False, C_INK),
            ("对方更优积木 → 我方写入底座 → 组装到终点", 15, True, C_WARM),
            ("。", 14, False, C_INK),
        ], [
            ("底座变好，终点跟着变好；不抢开放发现主语。", 13, False, C_MUTED),
        ],
    ], default_size=14)

    section_label(s, Inches(0.4), Inches(4.0), Inches(1.35), Inches(0.34), "口径提醒", fill_c=C_WARM_SOFT, line_c=C_WARM, text_c=C_WARM)
    rich(s, Inches(1.9), Inches(3.95), Inches(11.0), Inches(1.5), [
        [
            ("材料不写速度数字", 13, True, C_WARM),
            ("；口头可保守说：方法成型后，更短时间、更少额度做更大规模。", 13, False, C_INK),
        ], [
            ("检验床是 Ascend 上的 ML-KEM；", 13, False, C_INK),
            ("方法论服务一般预研，不绑定某一算法仓库名", 13, True, C_DEEP),
            ("。", 13, False, C_INK),
        ], [
            ("效果口径：", 13, False, C_DEEP),
            ("接近手调，相对更省", 13, True, C_WARM),
            ("——不作未对账的倍数宣称。", 13, False, C_INK),
        ],
    ], default_size=13)
    page_summary(s, "一条能力线，两段故事：先管住 AI，再有界往前推。", 10, 10)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"OK: {path}")
